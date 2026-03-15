#!/usr/bin/env python3
"""
Translate a .pptx presentation between languages (no API keys required).

Uses:
- python-pptx (read/write PowerPoint)
- deep-translator backends (Google + MyMemory fallback, no key)

Usage:
    python pptx_translate_ru.py input.pptx output_ru.pptx

Optional:
    python pptx_translate_ru.py input.pptx output_ru.pptx --src en --dest ru

Behavior:
- Retries transient translation failures
- Tries multiple translator backends for each chunk
- Translates paragraph/sentence text as a whole instead of word-by-word runs
- Splits only overly long paragraphs into smaller chunks
- Preserves original text for chunks that still fail
- Continues translating the rest of the deck instead of aborting on one bad sentence
"""

from __future__ import annotations

import argparse
import re
import sys
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, Dict, List, Optional, Tuple

from deep_translator import GoogleTranslator, MyMemoryTranslator
from pptx import Presentation


@dataclass
class TranslateConfig:
    src: str = "en"
    dest: str = "ru"
    retries: int = 3
    retry_delay_sec: float = 1.0
    chunk_limit: int = 4500
    max_reported_failures: int = 10


@dataclass
class TranslationFailure:
    location: str
    text: str
    error: str


@dataclass
class TranslationStats:
    translated_chunks: int = 0
    failed_chunks: int = 0
    failures: List[TranslationFailure] = field(default_factory=list)

    def add_success(self) -> None:
        self.translated_chunks += 1

    def add_failure(self, location: str, text: str, error: str, max_reported_failures: int) -> None:
        self.failed_chunks += 1
        if len(self.failures) < max_reported_failures:
            self.failures.append(
                TranslationFailure(
                    location=location,
                    text=text[:240],
                    error=error[:500],
                )
            )


@dataclass
class TranslatorBackend:
    name: str
    factory: Callable[[], object]


class SafeTranslator:
    def __init__(self, cfg: TranslateConfig):
        self.cfg = cfg
        self._cache: Dict[str, str] = {}
        self._backend_defs = [
            TranslatorBackend(
                name="google",
                factory=lambda: GoogleTranslator(source=cfg.src, target=cfg.dest),
            ),
            TranslatorBackend(
                name="mymemory",
                factory=lambda: MyMemoryTranslator(source=cfg.src, target=cfg.dest),
            ),
        ]
        self._backend_instances: Dict[str, object] = {}

    def translate_text(self, text: str, location: str, stats: TranslationStats) -> str:
        if not text or text.isspace():
            return text

        if text in self._cache:
            return self._cache[text]

        normalized_text = normalize_text_for_translation(text)
        translated = self._translate_normalized_text(normalized_text, location, stats)
        result = restore_text_shape(text, translated)
        self._cache[text] = result
        return result

    def _translate_normalized_text(self, text: str, location: str, stats: TranslationStats) -> str:
        if not text.strip():
            return text

        if len(text) <= self.cfg.chunk_limit:
            success, translated = self._try_translate(text, location, stats)
            return translated if success else text

        parts = split_long_text(text, self.cfg.chunk_limit)
        translated_parts: List[str] = []
        for part in parts:
            if not part.strip():
                translated_parts.append(part)
                continue

            success, translated = self._try_translate(part, location, stats)
            translated_parts.append(translated if success else part)

        return "".join(translated_parts)

    def _try_translate(self, text: str, location: str, stats: TranslationStats) -> Tuple[bool, str]:
        success, translated_or_error = self._translate_with_retries(text)
        if success:
            stats.add_success()
            return True, translated_or_error

        stats.add_failure(location, text, translated_or_error, self.cfg.max_reported_failures)
        return False, text

    def _translate_with_retries(self, text: str) -> Tuple[bool, str]:
        backend_errors: List[str] = []

        for backend in self._backend_defs:
            last_err: Optional[Exception] = None

            for attempt in range(1, self.cfg.retries + 1):
                try:
                    translator = self._get_backend_instance(backend)
                    translated = translator.translate(text)
                    if translated is None:
                        raise RuntimeError("Translator returned no text")
                    return True, translated
                except Exception as error:
                    last_err = error
                    if attempt < self.cfg.retries:
                        time.sleep(self.cfg.retry_delay_sec)

            backend_errors.append(
                f"{backend.name}: failed after {self.cfg.retries} attempts: {last_err}"
            )

        return False, " | ".join(backend_errors)

    def _get_backend_instance(self, backend: TranslatorBackend):
        existing = self._backend_instances.get(backend.name)
        if existing is not None:
            return existing

        instance = backend.factory()
        self._backend_instances[backend.name] = instance
        return instance


def normalize_text_for_translation(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def restore_text_shape(original_text: str, translated_text: str) -> str:
    leading_match = re.match(r"^\s*", original_text)
    trailing_match = re.search(r"\s*$", original_text)
    leading = leading_match.group(0) if leading_match else ""
    trailing = trailing_match.group(0) if trailing_match else ""
    return f"{leading}{translated_text.strip()}{trailing}"


def split_long_text(text: str, max_len: int) -> List[str]:
    sentences = re.split(r"(?<=[.!?;:])\s+|\n+", text)
    parts: List[str] = []
    current = ""

    for sentence in sentences:
        if not sentence:
            continue

        candidate = sentence if not current else f"{current} {sentence}"
        if len(candidate) <= max_len:
            current = candidate
            continue

        if current:
            parts.append(current)
            current = ""

        if len(sentence) <= max_len:
            current = sentence
            continue

        words = sentence.split(" ")
        segment = ""
        for word in words:
            candidate_word = word if not segment else f"{segment} {word}"
            if len(candidate_word) <= max_len:
                segment = candidate_word
            else:
                if segment:
                    parts.append(segment)
                if len(word) <= max_len:
                    segment = word
                else:
                    for index in range(0, len(word), max_len):
                        parts.append(word[index : index + max_len])
                    segment = ""

        if segment:
            current = segment

    if current:
        parts.append(current)

    return parts


def iter_text_containers(prs: Presentation):
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            base_location = f"slide {slide_index}, shape {shape_index}"

            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                yield shape.text_frame, base_location

            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                for row_index, row in enumerate(table.rows, start=1):
                    for cell_index, cell in enumerate(row.cells, start=1):
                        if cell.text_frame is not None:
                            yield cell.text_frame, f"{base_location}, table row {row_index}, cell {cell_index}"

            if shape.shape_type == 6:
                for subshape_index, subshape in enumerate(shape.shapes, start=1):
                    if hasattr(subshape, "has_text_frame") and subshape.has_text_frame:
                        yield subshape.text_frame, f"{base_location}, group item {subshape_index}"


def replace_paragraph_text(paragraph, translated_text: str) -> None:
    if paragraph.runs:
        paragraph.runs[0].text = translated_text
        for run in paragraph.runs[1:]:
            run.text = ""
        return

    paragraph.text = translated_text


def translate_text_frame(text_frame, location: str, translator: SafeTranslator, stats: TranslationStats) -> int:
    updated = 0
    for paragraph_index, paragraph in enumerate(text_frame.paragraphs, start=1):
        paragraph_location = f"{location}, paragraph {paragraph_index}"
        source_text = paragraph.text

        if not source_text or not source_text.strip():
            continue

        translated_text = translator.translate_text(source_text, paragraph_location, stats)
        if translated_text != source_text:
            replace_paragraph_text(paragraph, translated_text)
        updated += 1

    return updated


def translate_presentation(input_path: Path, output_path: Path, cfg: TranslateConfig) -> TranslationStats:
    presentation = Presentation(str(input_path))
    translator = SafeTranslator(cfg)
    stats = TranslationStats()

    changed = 0
    for text_frame, location in iter_text_containers(presentation):
        changed += translate_text_frame(text_frame, location, translator, stats)

    presentation.save(str(output_path))
    print(f"Done. Updated text chunks: {changed}")
    print(f"Successfully translated chunks: {stats.translated_chunks}")
    print(f"Failed chunks kept in original language: {stats.failed_chunks}")
    print(f"Saved: {output_path}")

    if stats.failures:
        print("Sample failures:", file=sys.stderr)
        for failure in stats.failures:
            print(
                f"- {failure.location}: {failure.error} | text={failure.text}",
                file=sys.stderr,
            )

    return stats


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Translate PPTX between languages without API keys")
    parser.add_argument("input", type=Path, help="Input .pptx file")
    parser.add_argument("output", type=Path, help="Output .pptx file")
    parser.add_argument("--src", default="en", help="Source language (default: en)")
    parser.add_argument("--dest", default="ru", help="Destination language (default: ru)")
    parser.add_argument("--retries", type=int, default=3, help="Retries per request")
    parser.add_argument("--retry-delay", type=float, default=1.0, help="Retry delay in seconds")
    return parser.parse_args()


def main() -> int:
    args = parse_args()

    if not args.input.exists():
        print(f"Input file not found: {args.input}", file=sys.stderr)
        return 1

    if args.input.suffix.lower() != ".pptx":
        print("Input must be a .pptx file", file=sys.stderr)
        return 1

    config = TranslateConfig(
        src=args.src,
        dest=args.dest,
        retries=max(1, args.retries),
        retry_delay_sec=max(0.0, args.retry_delay),
    )

    try:
        translate_presentation(args.input, args.output, config)
        return 0
    except Exception as error:
        print(f"Error: {error}", file=sys.stderr)
        return 2


if __name__ == "__main__":
    raise SystemExit(main())
